"""
PowerPoint file loader with text, table, and structure extraction.
Extracts content from slides including shapes, tables, and speaker notes.

Dependencies:
    pip install python-pptx
"""

import pandas as pd
from pathlib import Path
from typing import Tuple, Dict, Any, List, Optional
from dataclasses import dataclass, field
import re


@dataclass
class SlideContent:
    """Content extracted from a single slide."""
    slide_number: int
    title: Optional[str]
    body_text: List[str]
    tables: List[pd.DataFrame]
    notes: Optional[str]
    shapes_count: int


@dataclass  
class PPTXContent:
    """Structured content extracted from a PowerPoint file."""
    slides: List[SlideContent]
    all_text: str
    all_tables: List[pd.DataFrame]
    metadata: Dict[str, Any]
    warnings: List[str] = field(default_factory=list)


def extract_text_from_shape(shape) -> Optional[str]:
    """Extract text from a shape if it has a text frame."""
    if not shape.has_text_frame:
        return None
    
    paragraphs = []
    for paragraph in shape.text_frame.paragraphs:
        text = ''.join(run.text for run in paragraph.runs)
        if text.strip():
            paragraphs.append(text.strip())
    
    return '\n'.join(paragraphs) if paragraphs else None


def extract_table_from_shape(shape, slide_number: int) -> Optional[pd.DataFrame]:
    """Extract table data from a table shape."""
    if not shape.has_table:
        return None
    
    table = shape.table
    data = []
    
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            # Extract text from cell
            text = cell.text.strip() if cell.text else ""
            row_data.append(text)
        data.append(row_data)
    
    if not data or len(data) < 2:
        return None
    
    # First row as header
    try:
        df = pd.DataFrame(data[1:], columns=data[0])
        df = df.dropna(how='all').dropna(axis=1, how='all')
        
        if not df.empty:
            df.attrs['source_slide'] = slide_number
            return df
    except Exception:
        # If header row doesn't work, use numeric columns
        df = pd.DataFrame(data)
        df.attrs['source_slide'] = slide_number
        return df
    
    return None


def extract_slide_content(slide, slide_number: int) -> SlideContent:
    """Extract all content from a single slide."""
    title = None
    body_text = []
    tables = []
    shapes_count = 0
    
    for shape in slide.shapes:
        shapes_count += 1
        
        # Check for title
        if shape.is_placeholder:
            try:
                if shape.placeholder_format.type == 1:  # Title placeholder
                    text = extract_text_from_shape(shape)
                    if text:
                        title = text
                        continue
            except:
                pass
        
        # Check for table
        if shape.has_table:
            table_df = extract_table_from_shape(shape, slide_number)
            if table_df is not None:
                tables.append(table_df)
            continue
        
        # Extract text from other shapes
        text = extract_text_from_shape(shape)
        if text:
            body_text.append(text)
    
    # Extract speaker notes
    notes = None
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else None
        if notes_text and notes_text.strip():
            notes = notes_text.strip()
    
    return SlideContent(
        slide_number=slide_number,
        title=title,
        body_text=body_text,
        tables=tables,
        notes=notes,
        shapes_count=shapes_count
    )


def load_pptx(file_path: str) -> Tuple[PPTXContent, Dict[str, Any]]:
    """
    Load and extract content from a PowerPoint file.
    
    Args:
        file_path: Path to PPTX file
    
    Returns:
        Tuple of (PPTXContent, metadata dict)
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise ImportError("python-pptx required: pip install python-pptx")
    
    path = Path(file_path)
    
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")
    
    valid_extensions = ['.pptx', '.ppt']
    if path.suffix.lower() not in valid_extensions:
        raise ValueError(f"Expected PowerPoint file, got: {path.suffix}")
    
    warnings = []
    
    # Handle .ppt (old format)
    if path.suffix.lower() == '.ppt':
        warnings.append("Old .ppt format detected. Consider converting to .pptx for better extraction.")
    
    prs = Presentation(file_path)
    
    # Extract core properties if available
    pptx_metadata = {}
    try:
        core_props = prs.core_properties
        pptx_metadata = {
            'title': core_props.title,
            'author': core_props.author,
            'created': str(core_props.created) if core_props.created else None,
            'modified': str(core_props.modified) if core_props.modified else None,
            'subject': core_props.subject,
        }
    except:
        pass
    
    # Extract content from each slide
    slides_content = []
    all_text_parts = []
    all_tables = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_content = extract_slide_content(slide, i)
        slides_content.append(slide_content)
        
        # Aggregate text
        slide_text = f"\n--- Slide {i}"
        if slide_content.title:
            slide_text += f": {slide_content.title}"
        slide_text += " ---\n"
        
        if slide_content.title:
            slide_text += f"{slide_content.title}\n\n"
        
        if slide_content.body_text:
            slide_text += "\n".join(slide_content.body_text) + "\n"
        
        if slide_content.notes:
            slide_text += f"\n[Speaker Notes: {slide_content.notes}]\n"
        
        all_text_parts.append(slide_text)
        
        # Aggregate tables
        all_tables.extend(slide_content.tables)
    
    content = PPTXContent(
        slides=slides_content,
        all_text="\n".join(all_text_parts),
        all_tables=all_tables,
        metadata=pptx_metadata,
        warnings=warnings
    )
    
    metadata = {
        'source_file': path.name,
        'file_size_kb': round(path.stat().st_size / 1024, 2),
        'slides': len(slides_content),
        'tables_found': len(all_tables),
        'text_length': len(content.all_text),
        'has_speaker_notes': any(s.notes for s in slides_content),
        'pptx_metadata': pptx_metadata,
        'warnings': warnings,
    }
    
    return content, metadata


def pptx_to_dataframe(
    file_path: str,
    table_index: Optional[int] = None,
    combine_tables: bool = True
) -> Tuple[pd.DataFrame, Dict[str, Any]]:
    """
    Load PPTX and return extracted tables as DataFrame.
    
    Args:
        file_path: Path to PPTX file
        table_index: Specific table to return (0-indexed), or None for all/combined
        combine_tables: If True and multiple tables found, try to combine them
    
    Returns:
        Tuple of (DataFrame, metadata)
        DataFrame may be empty if no tables found
    """
    content, metadata = load_pptx(file_path)
    
    if not content.all_tables:
        metadata['note'] = "No tables found in presentation. Text content available."
        metadata['text_preview'] = content.all_text[:2000] if content.all_text else ""
        return pd.DataFrame(), metadata
    
    if table_index is not None:
        if table_index < len(content.all_tables):
            df = content.all_tables[table_index]
        else:
            raise IndexError(f"Table index {table_index} out of range. Found {len(content.all_tables)} tables.")
    elif combine_tables and len(content.all_tables) > 1:
        # Try to combine tables with same structure
        first_cols = set(content.all_tables[0].columns)
        compatible = [content.all_tables[0]]
        
        for table in content.all_tables[1:]:
            if set(table.columns) == first_cols:
                compatible.append(table)
        
        if len(compatible) == len(content.all_tables):
            df = pd.concat(content.all_tables, ignore_index=True)
        else:
            # Return largest table
            df = max(content.all_tables, key=len)
    else:
        df = content.all_tables[0]
    
    metadata['table_index'] = table_index
    metadata['rows_loaded'] = len(df)
    metadata['columns_loaded'] = len(df.columns)
    
    return df, metadata


def get_pptx_summary(file_path: str) -> Dict[str, Any]:
    """
    Quick summary of PowerPoint contents.
    Useful for previewing what's in a presentation.
    """
    content, metadata = load_pptx(file_path)
    
    # Slide summaries
    slide_summaries = []
    for slide in content.slides:
        slide_summaries.append({
            'slide_number': slide.slide_number,
            'title': slide.title,
            'text_items': len(slide.body_text),
            'tables': len(slide.tables),
            'has_notes': slide.notes is not None,
            'shapes': slide.shapes_count
        })
    
    # Table summaries
    table_summaries = []
    for i, table in enumerate(content.all_tables):
        table_summaries.append({
            'index': i,
            'rows': len(table),
            'columns': len(table.columns),
            'column_names': list(table.columns),
            'source_slide': table.attrs.get('source_slide', 'unknown')
        })
    
    return {
        'file_name': Path(file_path).name,
        'slides': len(content.slides),
        'slide_summaries': slide_summaries,
        'text_preview': content.all_text[:1000] + "..." if len(content.all_text) > 1000 else content.all_text,
        'text_length': len(content.all_text),
        'tables': table_summaries,
        'has_speaker_notes': metadata.get('has_speaker_notes', False),
        'pptx_metadata': content.metadata,
        'warnings': content.warnings
    }


def extract_metrics_from_pptx(file_path: str) -> Dict[str, Any]:
    """
    Extract business metrics and KPIs from a PowerPoint presentation.
    Looks for common patterns like revenue, percentages, growth rates.
    """
    content, _ = load_pptx(file_path)
    
    metrics = {
        'currencies': [],
        'percentages': [],
        'large_numbers': [],
        'kpi_candidates': []
    }
    
    text = content.all_text
    
    # Currency patterns
    currency_patterns = [
        r'(€|EUR)\s*([\d.,]+(?:\s*(?:Mio|Mrd|K|M|B))?)',
        r'(\$|USD)\s*([\d.,]+(?:\s*(?:K|M|B|million|billion))?)',
        r'([\d.,]+)\s*(€|EUR|USD|\$)',
    ]
    
    for pattern in currency_patterns:
        for match in re.finditer(pattern, text, re.IGNORECASE):
            # Try to find context (surrounding words)
            start = max(0, match.start() - 50)
            end = min(len(text), match.end() + 50)
            context = text[start:end].replace('\n', ' ')
            
            metrics['currencies'].append({
                'value': match.group(0),
                'context': context
            })
    
    # Percentage patterns
    pct_pattern = r'([\d.,]+)\s*(%|Prozent|percent)'
    for match in re.finditer(pct_pattern, text, re.IGNORECASE):
        start = max(0, match.start() - 50)
        end = min(len(text), match.end() + 50)
        context = text[start:end].replace('\n', ' ')
        
        metrics['percentages'].append({
            'value': match.group(1) + '%',
            'context': context
        })
    
    # Look for KPI-like patterns (word followed by colon and number)
    kpi_pattern = r'([A-Za-zÄÖÜäöü]+(?:\s+[A-Za-zÄÖÜäöü]+)?)\s*[:\-]\s*([\d.,]+(?:\s*(?:€|%|Mio|K|M))?)'
    for match in re.finditer(kpi_pattern, text):
        label = match.group(1).strip()
        value = match.group(2).strip()
        
        # Filter out noise
        if len(label) > 3 and len(value) > 0:
            metrics['kpi_candidates'].append({
                'label': label,
                'value': value
            })
    
    # Deduplicate
    seen = set()
    unique_kpis = []
    for kpi in metrics['kpi_candidates']:
        key = f"{kpi['label']}:{kpi['value']}"
        if key not in seen:
            seen.add(key)
            unique_kpis.append(kpi)
    metrics['kpi_candidates'] = unique_kpis
    
    return metrics

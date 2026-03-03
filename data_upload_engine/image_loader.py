"""
Image-based chart/graph data extractor for Local Analyst.
Extracts numeric data from screenshots, chart images, and embedded visuals.

OCR backends (in priority order):
  1. EasyOCR   — pip install easyocr  (pure Python, no system packages)
  2. pytesseract — pip install pytesseract + system tesseract-ocr (faster but needs OS install)
  3. No OCR    — graceful fallback, tells user what to pip install

Also uses:
  - OpenCV (optional) for image preprocessing
  - Pillow for image handling

Handles: bar charts, line charts, tables-as-images, KPI dashboards, screenshots.
"""

import pandas as pd
import numpy as np
from pathlib import Path
from typing import Tuple, Dict, Any, List, Optional
from dataclasses import dataclass, field
import re
import io


@dataclass
class ImageContent:
    """Structured content extracted from an image."""
    raw_text: str
    tables: List[pd.DataFrame]
    numbers: List[Dict[str, Any]]
    kpis: List[Dict[str, str]]
    metadata: Dict[str, Any]
    warnings: List[str] = field(default_factory=list)


# ── OCR backend detection ──

_ocr_backend = None  # 'easyocr', 'pytesseract', or None


def _detect_ocr_backend() -> Optional[str]:
    """Detect which OCR backend is available."""
    global _ocr_backend
    if _ocr_backend is not None:
        return _ocr_backend

    # Try EasyOCR first (pure pip, no system deps)
    try:
        import easyocr  # noqa: F401
        _ocr_backend = 'easyocr'
        return _ocr_backend
    except ImportError:
        pass

    # Try pytesseract (needs system tesseract binary)
    try:
        import pytesseract
        # Verify the system binary is actually available
        pytesseract.get_tesseract_version()
        _ocr_backend = 'pytesseract'
        return _ocr_backend
    except (ImportError, Exception):
        pass

    _ocr_backend = None
    return None


# Singleton EasyOCR reader (expensive to init, reuse across calls)
_easyocr_reader = None


def _get_easyocr_reader():
    global _easyocr_reader
    if _easyocr_reader is None:
        import easyocr
        _easyocr_reader = easyocr.Reader(['en', 'de'], gpu=False, verbose=False)
    return _easyocr_reader


def _preprocess_image(img_array):
    """Preprocess image for better OCR accuracy. Uses OpenCV if available, else Pillow."""
    try:
        import cv2

        if len(img_array.shape) == 3:
            gray = cv2.cvtColor(img_array, cv2.COLOR_BGR2GRAY)
        else:
            gray = img_array.copy()

        h, w = gray.shape[:2]
        if max(h, w) < 1000:
            scale = 2000 / max(h, w)
            gray = cv2.resize(gray, None, fx=scale, fy=scale, interpolation=cv2.INTER_CUBIC)

        binary = cv2.adaptiveThreshold(
            gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 15, 8
        )
        denoised = cv2.medianBlur(binary, 3)
        return denoised

    except ImportError:
        # Fallback: basic Pillow preprocessing
        from PIL import Image, ImageFilter, ImageOps
        if isinstance(img_array, np.ndarray):
            img = Image.fromarray(img_array)
        else:
            img = img_array

        img = ImageOps.grayscale(img)
        w, h = img.size
        if max(w, h) < 1000:
            scale = 2000 / max(w, h)
            img = img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
        img = img.filter(ImageFilter.SHARPEN)
        return np.array(img)


def _ocr_with_easyocr(img_array) -> str:
    """Run OCR using EasyOCR."""
    reader = _get_easyocr_reader()
    results = reader.readtext(img_array, detail=0, paragraph=True)
    return '\n'.join(results) if results else ''


def _ocr_with_pytesseract(img_array) -> str:
    """Run OCR using pytesseract."""
    import pytesseract
    results = []
    for psm in [6, 4, 3]:
        try:
            text = pytesseract.image_to_string(img_array, config=f'--psm {psm} --oem 3')
            if text.strip():
                results.append(text.strip())
        except Exception:
            continue
    return max(results, key=len) if results else ''


def _run_ocr(img_array) -> Tuple[str, str]:
    """
    Run OCR with best available backend.
    Returns (text, backend_used).
    """
    backend = _detect_ocr_backend()

    if backend == 'easyocr':
        return _ocr_with_easyocr(img_array), 'easyocr'
    elif backend == 'pytesseract':
        return _ocr_with_pytesseract(img_array), 'pytesseract'
    else:
        return '', 'none'


def _extract_text_ocr(image_path: str) -> Tuple[str, str]:
    """Run OCR on an image file. Returns (text, backend_used)."""
    from PIL import Image

    img = Image.open(image_path)
    img_array = np.array(img)

    # Preprocess
    preprocessed = _preprocess_image(img_array)

    return _run_ocr(preprocessed)


def _extract_text_from_bytes(image_bytes: bytes) -> str:
    """Run OCR on image bytes."""
    from PIL import Image

    img = Image.open(io.BytesIO(image_bytes))
    img_array = np.array(img)
    preprocessed = _preprocess_image(img_array)
    text, _ = _run_ocr(preprocessed)
    return text


def _parse_numbers_from_text(text: str) -> List[Dict[str, Any]]:
    """Extract labeled numbers from OCR text."""
    numbers = []

    # Pattern: label followed by number (with optional currency/unit)
    patterns = [
        # "Revenue: €1,234.56" or "Revenue 1234"
        r'([A-Za-zÄÖÜäöüß][A-Za-zÄÖÜäöüß\s\-/()]{2,40}?)\s*[:=\-–]\s*([€$£]?)\s*([\d.,]+)\s*(%|[KMBkmb](?:io)?|Mrd|million|billion)?',
        # "$1,234" or "€45.6M" standalone with nearby label
        r'([€$£])\s*([\d.,]+)\s*(%|[KMBkmb](?:io)?|Mrd|million|billion)?',
        # "45.6%" standalone
        r'([\d.,]+)\s*(%)',
    ]

    for pattern in patterns:
        for match in re.finditer(pattern, text, re.IGNORECASE):
            groups = match.groups()
            if len(groups) >= 3:
                label = groups[0].strip() if not groups[0].startswith(('€', '$', '£')) else None
                value_str = groups[-2] if len(groups) > 2 else groups[1]
                unit = groups[-1] if groups[-1] else None

                # Clean value
                try:
                    clean = value_str.replace(',', '').replace(' ', '')
                    val = float(clean)
                    numbers.append({
                        'label': label,
                        'value': val,
                        'raw': match.group(0).strip(),
                        'unit': unit
                    })
                except (ValueError, TypeError):
                    pass

    return numbers


def _parse_table_from_text(text: str) -> List[pd.DataFrame]:
    """Attempt to reconstruct tables from OCR text (tab/space-aligned)."""
    tables = []
    lines = [l for l in text.split('\n') if l.strip()]

    if len(lines) < 2:
        return tables

    # Strategy 1: detect tab-separated or multi-space-separated rows
    parsed_rows = []
    for line in lines:
        # Split by 2+ spaces or tabs
        cells = re.split(r'\t|  +', line.strip())
        cells = [c.strip() for c in cells if c.strip()]
        if len(cells) >= 2:
            parsed_rows.append(cells)

    if len(parsed_rows) >= 2:
        # Find most common column count
        col_counts = [len(r) for r in parsed_rows]
        from collections import Counter
        most_common_n = Counter(col_counts).most_common(1)[0][0]

        # Filter to rows with that column count
        consistent_rows = [r for r in parsed_rows if len(r) == most_common_n]

        if len(consistent_rows) >= 2:
            try:
                header = consistent_rows[0]
                data = consistent_rows[1:]
                df = pd.DataFrame(data, columns=header)

                # Try to convert numeric columns
                for col in df.columns:
                    try:
                        cleaned = df[col].str.replace(',', '').str.replace('€', '').str.replace('$', '').str.replace('%', '')
                        df[col] = pd.to_numeric(cleaned, errors='ignore')
                    except Exception:
                        pass

                if not df.empty:
                    tables.append(df)
            except Exception:
                pass

    return tables


def _extract_kpis(text: str) -> List[Dict[str, str]]:
    """Extract KPI-like label:value pairs from text."""
    kpis = []
    seen = set()

    patterns = [
        r'([A-Za-zÄÖÜäöüß][A-Za-zÄÖÜäöüß\s\-/()]{2,35}?)\s*[:=]\s*([€$£]?\s*[\d.,]+\s*(?:%|[KMBkmb](?:io)?|Mrd)?)',
        r'([A-Za-zÄÖÜäöüß][A-Za-zÄÖÜäöüß\s\-/()]{2,35}?)\s*\n\s*([€$£]?\s*[\d.,]+\s*(?:%|[KMBkmb](?:io)?|Mrd)?)',
    ]

    for pattern in patterns:
        for match in re.finditer(pattern, text, re.IGNORECASE):
            label = match.group(1).strip()
            value = match.group(2).strip()

            # Skip noise
            if len(label) < 3 or label.lower() in ('page', 'slide', 'date', 'time', 'figure', 'chart'):
                continue

            key = f"{label}:{value}"
            if key not in seen:
                seen.add(key)
                kpis.append({'label': label, 'value': value})

    return kpis


def load_image(file_path: str) -> Tuple[ImageContent, Dict[str, Any]]:
    """
    Load an image file and extract data using OCR.

    Args:
        file_path: Path to image file (JPEG, PNG, BMP, TIFF, WEBP)

    Returns:
        Tuple of (ImageContent, metadata dict)
    """
    from PIL import Image

    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    valid_exts = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp'}
    if path.suffix.lower() not in valid_exts:
        raise ValueError(f"Unsupported image format: {path.suffix}")

    warnings = []

    # Get image info
    img = Image.open(file_path)
    width, height = img.size
    img_format = img.format or path.suffix.upper().strip('.')

    if max(width, height) < 200:
        warnings.append("Image is very small — OCR accuracy may be low.")

    # OCR
    backend = _detect_ocr_backend()
    if backend is None:
        warnings.append(
            "No OCR backend installed. Install one with:\n"
            "  pip install easyocr          (recommended — pure Python, no system packages)\n"
            "  pip install pytesseract       (faster, but also needs: apt install tesseract-ocr)"
        )
        raw_text = ""
        ocr_backend_used = 'none'
    else:
        try:
            raw_text, ocr_backend_used = _extract_text_ocr(file_path)
        except Exception as e:
            raw_text = ""
            ocr_backend_used = 'error'
            warnings.append(f"OCR failed: {e}")

    if not raw_text:
        warnings.append("No text could be extracted from the image. "
                         "This may be a chart with minimal text labels.")

    # Parse structured data from text
    tables = _parse_table_from_text(raw_text)
    numbers = _parse_numbers_from_text(raw_text)
    kpis = _extract_kpis(raw_text)

    content = ImageContent(
        raw_text=raw_text,
        tables=tables,
        numbers=numbers,
        kpis=kpis,
        metadata={'format': img_format, 'width': width, 'height': height},
        warnings=warnings
    )

    metadata = {
        'source_file': path.name,
        'file_size_kb': round(path.stat().st_size / 1024, 2),
        'image_format': img_format,
        'dimensions': f"{width}x{height}",
        'ocr_backend': ocr_backend_used,
        'text_extracted': len(raw_text),
        'tables_found': len(tables),
        'numbers_found': len(numbers),
        'kpis_found': len(kpis),
        'warnings': warnings,
    }

    return content, metadata


def image_to_dataframe(file_path: str) -> Tuple[pd.DataFrame, Dict[str, Any]]:
    """
    Load image and return best-effort DataFrame.

    If a table is detected via OCR, returns it.
    Otherwise, returns a DataFrame of extracted KPIs/numbers.
    """
    content, metadata = load_image(file_path)

    if content.tables:
        df = content.tables[0]
        metadata['extraction_type'] = 'table_ocr'
        return df, metadata

    # Build a DataFrame from extracted numbers/KPIs
    if content.kpis:
        df = pd.DataFrame(content.kpis)
        metadata['extraction_type'] = 'kpi_extraction'
        return df, metadata

    if content.numbers:
        rows = [{'label': n.get('label', ''), 'value': n['value'],
                 'unit': n.get('unit', ''), 'raw': n.get('raw', '')}
                for n in content.numbers]
        df = pd.DataFrame(rows)
        metadata['extraction_type'] = 'number_extraction'
        return df, metadata

    # Last resort: return raw text lines as single-column DataFrame
    if content.raw_text:
        lines = [l.strip() for l in content.raw_text.split('\n') if l.strip()]
        df = pd.DataFrame({'extracted_text': lines})
        metadata['extraction_type'] = 'raw_text'
        return df, metadata

    metadata['extraction_type'] = 'empty'
    return pd.DataFrame(), metadata


def extract_images_from_pptx(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract embedded images from a PowerPoint file and OCR each one.

    Returns list of dicts with: slide_number, image_bytes, ocr_text, numbers, kpis
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    extracted = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    blob = image.blob
                    ext = image.content_type.split('/')[-1]  # e.g. 'png', 'jpeg'

                    # OCR the image bytes
                    try:
                        text = _extract_text_from_bytes(blob)
                    except Exception:
                        text = ""

                    numbers = _parse_numbers_from_text(text) if text else []
                    kpis = _extract_kpis(text) if text else []

                    extracted.append({
                        'slide_number': slide_idx,
                        'image_format': ext,
                        'image_bytes': blob,
                        'size_kb': round(len(blob) / 1024, 1),
                        'ocr_text': text,
                        'numbers': numbers,
                        'kpis': kpis,
                    })
                except Exception:
                    continue

            # Also check for chart objects (EMF/WMF embedded charts)
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                try:
                    chart = shape.chart
                    chart_data = _extract_chart_object_data(chart)
                    if chart_data is not None and not chart_data.empty:
                        extracted.append({
                            'slide_number': slide_idx,
                            'image_format': 'chart_object',
                            'image_bytes': None,
                            'size_kb': 0,
                            'ocr_text': '',
                            'numbers': [],
                            'kpis': [],
                            'chart_data': chart_data,
                        })
                except Exception:
                    continue

    return extracted


def _extract_chart_object_data(chart) -> Optional[pd.DataFrame]:
    """Extract data directly from a python-pptx Chart object."""
    try:
        plot = chart.plots[0]
        categories = [str(c) for c in chart.plots[0].categories]

        data = {'Category': categories}
        for series in plot.series:
            name = series.name if series.name else f"Series {len(data)}"
            data[name] = list(series.values)

        df = pd.DataFrame(data)
        if not df.empty:
            return df
    except Exception:
        pass
    return None


def extract_all_image_data_from_pptx(file_path: str) -> Tuple[List[pd.DataFrame], Dict[str, Any]]:
    """
    Extract ALL data from images and charts in a PPTX file.
    Returns combined DataFrames and summary metadata.
    """
    image_items = extract_images_from_pptx(file_path)

    tables = []
    summary = {
        'images_found': len(image_items),
        'charts_with_data': 0,
        'images_with_text': 0,
        'total_kpis': 0,
        'total_numbers': 0,
        'details': []
    }

    for item in image_items:
        detail = {
            'slide': item['slide_number'],
            'format': item['image_format'],
            'size_kb': item['size_kb']
        }

        # Native chart data (best quality)
        if item.get('chart_data') is not None:
            tables.append(item['chart_data'])
            summary['charts_with_data'] += 1
            detail['type'] = 'native_chart'
            detail['rows'] = len(item['chart_data'])

        # OCR-extracted text
        elif item.get('ocr_text'):
            summary['images_with_text'] += 1
            ocr_tables = _parse_table_from_text(item['ocr_text'])
            if ocr_tables:
                tables.extend(ocr_tables)
                detail['type'] = 'ocr_table'
                detail['rows'] = sum(len(t) for t in ocr_tables)
            elif item.get('kpis'):
                kpi_df = pd.DataFrame(item['kpis'])
                kpi_df['source_slide'] = item['slide_number']
                tables.append(kpi_df)
                detail['type'] = 'ocr_kpis'
                detail['kpis'] = len(item['kpis'])

        summary['total_kpis'] += len(item.get('kpis', []))
        summary['total_numbers'] += len(item.get('numbers', []))
        summary['details'].append(detail)

    return tables, summary

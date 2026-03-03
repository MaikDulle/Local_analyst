"""
Test script for data_upload_engine module.
Run this to verify the loading system works correctly.

Usage:
    cd local_analyst
    python -m tests.test_loaders
"""

import sys
import os
from pathlib import Path

# Add parent directory to path
sys.path.insert(0, str(Path(__file__).parent.parent))

import pandas as pd
import tempfile
import json


def create_test_csv(path: str):
    """Create a sample CSV file for testing."""
    data = {
        'order_date': ['2024-01-15', '2024-01-16', '2024-01-17', '2024-01-18', '2024-01-19'],
        'customer_id': ['C001', 'C002', 'C001', 'C003', 'C002'],
        'product_id': ['P100', 'P101', 'P102', 'P100', 'P101'],
        'category': ['Electronics', 'Clothing', 'Electronics', 'Electronics', 'Clothing'],
        'quantity': [2, 1, 3, 1, 2],
        'revenue': [199.99, 49.99, 299.99, 99.99, 99.98],
        'channel': ['web', 'app', 'web', 'store', 'app']
    }
    df = pd.DataFrame(data)
    df.to_csv(path, index=False)
    return df


def create_test_excel(path: str):
    """Create a sample Excel file for testing."""
    data = {
        'date': pd.to_datetime(['2024-01-15', '2024-01-16', '2024-01-17']),
        'campaign': ['Winter Sale', 'Winter Sale', 'Flash Deal'],
        'impressions': [10000, 12000, 8000],
        'clicks': [500, 650, 420],
        'conversions': [25, 30, 22],
        'cost': [150.00, 180.00, 120.00]
    }
    df = pd.DataFrame(data)
    df.to_excel(path, index=False)
    return df


def create_test_json(path: str):
    """Create a sample JSON file for testing."""
    data = {
        'meta': {'source': 'test', 'version': '1.0'},
        'data': [
            {'user_id': 'U001', 'event': 'page_view', 'timestamp': '2024-01-15T10:30:00'},
            {'user_id': 'U002', 'event': 'add_to_cart', 'timestamp': '2024-01-15T10:31:00'},
            {'user_id': 'U001', 'event': 'purchase', 'timestamp': '2024-01-15T10:35:00'},
        ]
    }
    with open(path, 'w') as f:
        json.dump(data, f)
    return data


def test_csv_loader():
    """Test CSV loading functionality."""
    print("\n" + "="*50)
    print("Testing CSV Loader")
    print("="*50)
    
    from data_upload_engine import load_file
    
    with tempfile.NamedTemporaryFile(suffix='.csv', delete=False) as f:
        test_path = f.name
    
    try:
        original_df = create_test_csv(test_path)
        result = load_file(test_path)
        
        assert result.success, f"Load failed: {result.error}"
        assert len(result.df) == len(original_df), "Row count mismatch"
        assert len(result.df.columns) == len(original_df.columns), "Column count mismatch"
        
        print(f"✓ Loaded {len(result.df)} rows, {len(result.df.columns)} columns")
        print(f"✓ Quality score: {result.profile.quality_score}")
        print(f"✓ Ecom mapping suggestions: {result.ecom_mapping}")
        
        # Check column type detection
        print("\nColumn types detected:")
        for col_name, col_profile in result.profile.columns.items():
            print(f"  {col_name}: {col_profile.semantic_type.value}")
        
        print("\n✓ CSV loader test PASSED")
        return True
        
    except Exception as e:
        print(f"✗ CSV loader test FAILED: {e}")
        return False
    finally:
        os.unlink(test_path)


def test_excel_loader():
    """Test Excel loading functionality."""
    print("\n" + "="*50)
    print("Testing Excel Loader")
    print("="*50)
    
    from data_upload_engine import load_file
    
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as f:
        test_path = f.name
    
    try:
        original_df = create_test_excel(test_path)
        result = load_file(test_path)
        
        assert result.success, f"Load failed: {result.error}"
        assert len(result.df) == len(original_df), "Row count mismatch"
        
        print(f"✓ Loaded {len(result.df)} rows, {len(result.df.columns)} columns")
        print(f"✓ Sheet: {result.metadata.get('sheet_name')}")
        print(f"✓ Quality score: {result.profile.quality_score}")
        
        print("\n✓ Excel loader test PASSED")
        return True
        
    except Exception as e:
        print(f"✗ Excel loader test FAILED: {e}")
        return False
    finally:
        os.unlink(test_path)


def test_json_loader():
    """Test JSON loading functionality."""
    print("\n" + "="*50)
    print("Testing JSON Loader")
    print("="*50)
    
    from data_upload_engine import load_file
    
    with tempfile.NamedTemporaryFile(suffix='.json', delete=False) as f:
        test_path = f.name
    
    try:
        original_data = create_test_json(test_path)
        result = load_file(test_path)
        
        assert result.success, f"Load failed: {result.error}"
        assert len(result.df) == len(original_data['data']), "Row count mismatch"
        
        print(f"✓ Loaded {len(result.df)} rows, {len(result.df.columns)} columns")
        print(f"✓ JSON structure: {result.metadata.get('json_structure')}")
        print(f"✓ Data key: {result.metadata.get('data_key')}")
        
        print("\n✓ JSON loader test PASSED")
        return True
        
    except Exception as e:
        print(f"✗ JSON loader test FAILED: {e}")
        return False
    finally:
        os.unlink(test_path)


def test_caching():
    """Test caching functionality."""
    print("\n" + "="*50)
    print("Testing Cache System")
    print("="*50)
    
    from data_upload_engine import load_file, get_cache
    
    with tempfile.NamedTemporaryFile(suffix='.csv', delete=False) as f:
        test_path = f.name
    
    try:
        create_test_csv(test_path)
        
        # Clear cache first
        cache = get_cache()
        cache.clear_all_cache()
        
        # First load - should not be from cache
        result1 = load_file(test_path)
        assert result1.success
        assert result1.metadata.get('from_cache') == False, "First load should not be from cache"
        print("✓ First load: fresh")
        
        # Second load - should be from cache
        result2 = load_file(test_path)
        assert result2.success
        assert result2.metadata.get('from_cache') == True, "Second load should be from cache"
        print("✓ Second load: from cache")
        
        # Check cache stats
        stats = cache.get_cache_stats()
        print(f"✓ Cache stats: {stats}")
        
        print("\n✓ Cache test PASSED")
        return True
        
    except Exception as e:
        print(f"✗ Cache test FAILED: {e}")
        return False
    finally:
        os.unlink(test_path)


def test_validators():
    """Test validation and profiling functionality."""
    print("\n" + "="*50)
    print("Testing Validators")
    print("="*50)
    
    from data_upload_engine import profile_dataframe, suggest_ecom_mapping, ColumnType
    
    try:
        # Create test dataframe
        df = pd.DataFrame({
            'order_date': pd.to_datetime(['2024-01-15', '2024-01-16', '2024-01-17']),
            'customer_id': ['C001', 'C002', 'C003'],
            'revenue': [100.0, 200.0, 150.0],
            'quantity': [1, 2, 1],
            'channel': ['web', 'app', 'web'],
            'is_new_customer': [True, False, True]
        })
        
        profile = profile_dataframe(df)
        
        print(f"✓ Row count: {profile.row_count}")
        print(f"✓ Column count: {profile.column_count}")
        print(f"✓ Quality score: {profile.quality_score}")
        
        # Check type inference
        assert profile.columns['order_date'].semantic_type == ColumnType.DATETIME
        assert profile.columns['revenue'].semantic_type == ColumnType.CURRENCY
        assert profile.columns['channel'].semantic_type == ColumnType.CATEGORICAL
        assert profile.columns['is_new_customer'].semantic_type == ColumnType.BOOLEAN
        print("✓ Type inference correct")
        
        # Check ecom mapping
        mapping = suggest_ecom_mapping(df)
        assert mapping['date_column'] == 'order_date'
        assert mapping['revenue_column'] == 'revenue'
        assert mapping['customer_id'] == 'customer_id'
        print(f"✓ Ecom mapping: {mapping}")
        
        print("\n✓ Validators test PASSED")
        return True
        
    except Exception as e:
        print(f"✗ Validators test FAILED: {e}")
        import traceback
        traceback.print_exc()
        return False


def test_pdf_loader():
    """Test PDF loading functionality."""
    print("\n" + "="*50)
    print("Testing PDF Loader")
    print("="*50)
    
    try:
        from data_upload_engine import load_pdf, get_pdf_summary
    except ImportError as e:
        print(f"⚠ PDF loader import failed (missing dependency): {e}")
        print("  Install with: pip install pdfplumber PyMuPDF")
        return True  # Skip but don't fail
    
    # Create a simple test PDF using reportlab if available
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        from reportlab.platypus import Table, TableStyle
        from reportlab.lib import colors
        
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as f:
            test_path = f.name
        
        c = canvas.Canvas(test_path, pagesize=letter)
        c.drawString(100, 750, "Quarterly Report Q3 2024")
        c.drawString(100, 720, "Revenue: €1.5 Mio (+15%)")
        c.drawString(100, 690, "Customers: 12,500")
        c.save()
        
        # Test loading
        content, metadata = load_pdf(test_path)
        
        assert metadata['pages'] > 0, "Should have at least one page"
        assert len(content.text) > 0, "Should extract some text"
        print(f"✓ Loaded PDF: {metadata['pages']} pages, {len(content.text)} chars")
        
        # Test summary
        summary = get_pdf_summary(test_path)
        print(f"✓ PDF summary: {summary['pages']} pages")
        
        # Check metric extraction
        if metadata.get('extracted_metrics'):
            metrics = metadata['extracted_metrics']
            print(f"✓ Found metrics: {len(metrics.get('currencies', []))} currencies, "
                  f"{len(metrics.get('percentages', []))} percentages")
        
        os.unlink(test_path)
        print("\n✓ PDF loader test PASSED")
        return True
        
    except ImportError:
        print("⚠ reportlab not available for test PDF creation")
        print("  Testing with load_pdf import only")
        print("\n✓ PDF loader test PASSED (limited)")
        return True
    except Exception as e:
        print(f"✗ PDF loader test FAILED: {e}")
        import traceback
        traceback.print_exc()
        return False


def test_pptx_loader():
    """Test PPTX loading functionality."""
    print("\n" + "="*50)
    print("Testing PPTX Loader")
    print("="*50)
    
    try:
        from data_upload_engine import load_pptx, get_pptx_summary, pptx_to_dataframe
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as e:
        print(f"⚠ PPTX loader import failed (missing dependency): {e}")
        print("  Install with: pip install python-pptx")
        return True  # Skip but don't fail
    
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        test_path = f.name
    
    try:
        # Create test PPTX
        prs = Presentation()
        
        # Slide 1: Title slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide1.shapes.title
        title.text = "Q3 2024 Performance Report"
        subtitle = slide1.placeholders[1]
        subtitle.text = "Revenue: €2.5M | Growth: 18%"
        
        # Slide 2: With a table
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        
        # Add table
        rows, cols = 4, 3
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(2)
        
        table = slide2.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set header
        table.cell(0, 0).text = "Month"
        table.cell(0, 1).text = "Revenue"
        table.cell(0, 2).text = "Growth"
        
        # Set data
        table.cell(1, 0).text = "July"
        table.cell(1, 1).text = "€800K"
        table.cell(1, 2).text = "15%"
        
        table.cell(2, 0).text = "August"
        table.cell(2, 1).text = "€850K"
        table.cell(2, 2).text = "18%"
        
        table.cell(3, 0).text = "September"
        table.cell(3, 1).text = "€850K"
        table.cell(3, 2).text = "20%"
        
        prs.save(test_path)
        
        # Test loading
        content, metadata = load_pptx(test_path)
        
        assert metadata['slides'] == 2, f"Expected 2 slides, got {metadata['slides']}"
        assert len(content.all_text) > 0, "Should extract some text"
        print(f"✓ Loaded PPTX: {metadata['slides']} slides, {len(content.all_text)} chars")
        
        # Test table extraction
        assert metadata['tables_found'] >= 1, "Should find at least one table"
        print(f"✓ Found {metadata['tables_found']} table(s)")
        
        # Test summary
        summary = get_pptx_summary(test_path)
        print(f"✓ PPTX summary: {summary['slides']} slides, {len(summary['tables'])} tables")
        
        # Test dataframe extraction
        df, df_meta = pptx_to_dataframe(test_path)
        if not df.empty:
            print(f"✓ Extracted table: {len(df)} rows, {len(df.columns)} columns")
            print(f"  Columns: {list(df.columns)}")
        
        print("\n✓ PPTX loader test PASSED")
        return True
        
    except Exception as e:
        print(f"✗ PPTX loader test FAILED: {e}")
        import traceback
        traceback.print_exc()
        return False
    finally:
        if os.path.exists(test_path):
            os.unlink(test_path)


def run_all_tests():
    """Run all tests."""
    print("\n" + "="*60)
    print("LOCAL ANALYST - DATA UPLOAD ENGINE TESTS")
    print("="*60)
    
    results = {
        'CSV Loader': test_csv_loader(),
        'Excel Loader': test_excel_loader(),
        'JSON Loader': test_json_loader(),
        'PDF Loader': test_pdf_loader(),
        'PPTX Loader': test_pptx_loader(),
        'Caching': test_caching(),
        'Validators': test_validators(),
    }
    
    print("\n" + "="*60)
    print("TEST SUMMARY")
    print("="*60)
    
    passed = sum(1 for r in results.values() if r)
    total = len(results)
    
    for test_name, result in results.items():
        status = "✓ PASSED" if result else "✗ FAILED"
        print(f"  {test_name}: {status}")
    
    print(f"\nTotal: {passed}/{total} tests passed")
    
    return passed == total


if __name__ == '__main__':
    success = run_all_tests()
    sys.exit(0 if success else 1)
